#!/usr/bin/env python3
"""
Movie Recommendation API - PowerPoint Presentation Generator
Creates a beautiful, professional presentation with all slides and content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

class MovieAPIPresentation:
    def __init__(self):
        self.prs = Presentation()
        self.setup_slide_size()
        
    def setup_slide_size(self):
        """Set up 16:9 aspect ratio"""
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self):
        """Slide 1: Title Slide"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 15, 26)  # Dark blue
        
        # Title
        title = slide.shapes.title
        title.text = "Movie Recommendation API"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = "Database Design & Implementation"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
        
        # Add project details
        self.add_text_box(slide, "Student: [Your Name]", 1, 5, 2, 0.5)
        self.add_text_box(slide, "Project: Shimy Movies", 1, 5.5, 2, 0.5)
        self.add_text_box(slide, "Technology: Django REST Framework, PostgreSQL, Redis", 1, 6, 2, 0.5)
        self.add_text_box(slide, "Date: [Current Date]", 1, 6.5, 2, 0.5)
        
        # Add decorative elements
        self.add_decorative_elements(slide)
    
    def add_project_overview_slide(self):
        """Slide 2: Project Overview"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Background
        self.set_dark_background(slide)
        
        # Title
        title = slide.shapes.title
        title.text = "Project Overview"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Content
        content = slide.placeholders[1]
        content.text = "What is Shimy Movies?"
        content.text_frame.paragraphs[0].font.size = Pt(20)
        content.text_frame.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)
        content.text_frame.paragraphs[0].font.bold = True
        
        # Add bullet points
        self.add_bullet_points(slide, [
            "Full-stack movie discovery platform",
            "TMDB API integration for comprehensive movie data",
            "Personalized recommendations based on user preferences",
            "Modern React frontend with Next.js",
            "RESTful API backend with Django"
        ], 1, 2.5, 5, 3)
        
        # Key Features section
        self.add_text_box(slide, "Key Features:", 7, 2.5, 2, 0.5)
        self.add_bullet_points(slide, [
            "User authentication & profiles",
            "Movie browsing & search",
            "Favorites & watchlist management",
            "Movie ratings & reviews",
            "Personalized recommendations"
        ], 7, 3, 2, 2.5)
    
    def add_technology_stack_slide(self):
        """Slide 3: Technology Stack"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Technology Stack"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Backend Technologies
        self.add_text_box(slide, "Backend Technologies:", 1, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "Django 4.2.7 - Web framework",
            "Django REST Framework - API development",
            "PostgreSQL - Primary database",
            "Redis - Caching & session storage",
            "JWT Authentication - Secure user sessions"
        ], 1, 2.5, 3, 2.5)
        
        # Frontend Technologies
        self.add_text_box(slide, "Frontend Technologies:", 5, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "Next.js - React framework",
            "TypeScript - Type safety",
            "Styled Components - Styling",
            "Vercel - Deployment platform"
        ], 5, 2.5, 3, 2)
        
        # External APIs
        self.add_text_box(slide, "External APIs:", 9, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "TMDB API - Movie data source",
            "YouTube API - Trailer integration"
        ], 9, 2.5, 3, 1)
    
    def set_dark_background(self, slide):
        """Set dark background for slide"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 15, 26)
    
    def add_text_box(self, slide, text, left, top, width, height):
        """Add a text box to the slide"""
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    def add_bullet_points(self, slide, points, left, top, width, height):
        """Add bullet points to the slide"""
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.clear()
        
        for i, point in enumerate(points):
            p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.alignment = PP_ALIGN.LEFT
    
    def add_decorative_elements(self, slide):
        """Add decorative elements to the title slide"""
        # Add a decorative rectangle
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1), Inches(12), Inches(0.1)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(52, 152, 219)
        rect.line.fill.background()
    
    def create_presentation(self):
        """Create the complete presentation"""
        print("🎬 Creating Movie Recommendation API Presentation...")
        
        # Add slides
        self.add_title_slide()
        self.add_project_overview_slide()
        self.add_technology_stack_slide()
        
        # Save the presentation
        output_path = "Movie_Recommendation_API_Presentation.pptx"
        self.prs.save(output_path)
        print(f"✅ Presentation saved as: {output_path}")
        print(f"📊 Total slides created: {len(self.prs.slides)}")
        
        return output_path

def main():
    """Main function to create the presentation"""
    try:
        # Create presentation
        presentation = MovieAPIPresentation()
        output_file = presentation.create_presentation()
        
        print("\n🎉 Presentation created successfully!")
        print(f"📁 File location: {os.path.abspath(output_file)}")
        print("\n📋 Next steps:")
        print("1. Open the PowerPoint file")
        print("2. Review and customize content")
        print("3. Add your personal information")
        print("4. Add screenshots or diagrams")
        print("5. Practice your presentation")
        
    except Exception as e:
        print(f"❌ Error creating presentation: {str(e)}")
        print("Make sure you have python-pptx installed:")
        print("pip install python-pptx")

if __name__ == "__main__":
    main() 