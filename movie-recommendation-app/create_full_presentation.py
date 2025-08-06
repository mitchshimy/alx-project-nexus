#!/usr/bin/env python3
"""
Movie Recommendation API - Complete PowerPoint Presentation Generator
Creates a beautiful, professional presentation with all 18 slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

class MovieAPIPresentation:
    def __init__(self):
        self.prs = Presentation()
        self.setup_slide_size()
        
    def setup_slide_size(self):
        """Set up 16:9 aspect ratio"""
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self):
        """Slide 1: Title Slide"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 15, 26)
        
        # Title
        title = slide.shapes.title
        title.text = "Movie Recommendation API"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = "Database Design & Implementation"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
        
        # Project details
        self.add_text_box(slide, "Student: [Your Name]", 1, 5, 2, 0.5)
        self.add_text_box(slide, "Project: Shimy Movies", 1, 5.5, 2, 0.5)
        self.add_text_box(slide, "Technology: Django REST Framework, PostgreSQL, Redis", 1, 6, 2, 0.5)
        self.add_text_box(slide, "Date: [Current Date]", 1, 6.5, 2, 0.5)
        
        # Decorative elements
        self.add_decorative_elements(slide)
    
    def add_project_overview_slide(self):
        """Slide 2: Project Overview"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Project Overview"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        content = slide.placeholders[1]
        content.text = "What is Shimy Movies?"
        content.text_frame.paragraphs[0].font.size = Pt(20)
        content.text_frame.paragraphs[0].font.color.rgb = RGBColor(52, 152, 219)
        content.text_frame.paragraphs[0].font.bold = True
        
        self.add_bullet_points(slide, [
            "Full-stack movie discovery platform",
            "TMDB API integration for comprehensive movie data",
            "Personalized recommendations based on user preferences",
            "Modern React frontend with Next.js",
            "RESTful API backend with Django"
        ], 1, 2.5, 5, 3)
        
        self.add_text_box(slide, "Key Features:", 7, 2.5, 2, 0.5)
        self.add_bullet_points(slide, [
            "User authentication & profiles",
            "Movie browsing & search",
            "Favorites & watchlist management",
            "Movie ratings & reviews",
            "Personalized recommendations"
        ], 7, 3, 2, 2.5)
    
    def add_technology_stack_slide(self):
        """Slide 3: Technology Stack"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Technology Stack"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Backend Technologies:", 1, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "Django 4.2.7 - Web framework",
            "Django REST Framework - API development",
            "PostgreSQL - Primary database",
            "Redis - Caching & session storage",
            "JWT Authentication - Secure user sessions"
        ], 1, 2.5, 3, 2.5)
        
        self.add_text_box(slide, "Frontend Technologies:", 5, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "Next.js - React framework",
            "TypeScript - Type safety",
            "Styled Components - Styling",
            "Vercel - Deployment platform"
        ], 5, 2.5, 3, 2)
        
        self.add_text_box(slide, "External APIs:", 9, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "TMDB API - Movie data source",
            "YouTube API - Trailer integration"
        ], 9, 2.5, 3, 1)
    
    def add_database_architecture_slide(self):
        """Slide 4: Database Architecture Overview"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Database Architecture Overview"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Core Design Principles:", 1, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "Normalized database design (3NF)",
            "Scalable architecture for large datasets",
            "Efficient query patterns for user preferences",
            "Flexible JSON fields for complex data",
            "Referential integrity with foreign keys"
        ], 1, 2.5, 3, 2.5)
        
        self.add_text_box(slide, "Database Entities:", 5, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "User & UserProfile",
            "Movie (with TMDB integration)",
            "Favorite (junction table)",
            "Watchlist (junction table)",
            "MovieRating (with reviews)"
        ], 5, 2.5, 3, 2.5)
    
    def add_erd_slide(self):
        """Slide 5: Entity Relationship Diagram"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Entity Relationship Diagram"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Database Schema:", 1, 2, 4, 0.5)
        
        erd_text = """User (1) ←→ (1) UserProfile
   ↓
   ↓ (1:N)
   ↓
Favorite ←→ (N:1) Movie
Watchlist ←→ (N:1) Movie
MovieRating ←→ (N:1) Movie"""
        
        self.add_text_box(slide, erd_text, 1, 2.5, 4, 2)
        
        self.add_text_box(slide, "Key Relationships:", 6, 2, 4, 0.5)
        self.add_bullet_points(slide, [
            "One-to-One: User ↔ UserProfile",
            "One-to-Many: User → Favorites/Watchlist/Ratings",
            "Many-to-Many: User ↔ Movie (through junction tables)"
        ], 6, 2.5, 4, 1.5)
    
    def add_user_management_slide(self):
        """Slide 6: User Management System"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "User Management System"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "User Entity:", 1, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "Extends Django's AbstractUser",
            "Email-based authentication",
            "JWT token management",
            "Profile extension capabilities"
        ], 1, 2.5, 3, 2)
        
        self.add_text_box(slide, "UserProfile Entity:", 5, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "Extended user preferences",
            "Favorite genres (JSON)",
            "Notification settings (JSON)",
            "Bio and personal information"
        ], 5, 2.5, 3, 2)
        
        self.add_text_box(slide, "Security Features:", 9, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "Password hashing",
            "Token expiration",
            "Rate limiting",
            "Input validation"
        ], 9, 2.5, 3, 2)
    
    def add_movie_data_model_slide(self):
        """Slide 7: Movie Data Model"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Movie Data Model"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Movie Entity Features:", 1, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "TMDB Integration: tmdb_id for external API sync",
            "Comprehensive Data: Title, overview, metadata",
            "Media Support: Movies and TV shows",
            "Rich Information: Budget, revenue, cast, crew",
            "Performance Metrics: Ratings, popularity, vote counts"
        ], 1, 2.5, 3, 2.5)
        
        self.add_text_box(slide, "Data Sources:", 5, 2, 3, 0.5)
        self.add_bullet_points(slide, [
            "TMDB API for movie information",
            "Real-time data synchronization",
            "Cached responses for performance",
            "JSON fields for flexible metadata"
        ], 5, 2.5, 3, 2)
    
    def add_user_interaction_slide(self):
        """Slide 8: User Interaction Models"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "User Interaction Models"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Three Core Interaction Types:", 1, 2, 4, 0.5)
        
        self.add_text_box(slide, "1. Favorites System", 1, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Junction table: User ↔ Movie",
            "Unique constraints prevent duplicates",
            "Timestamp tracking for analytics"
        ], 1, 3, 4, 1.5)
        
        self.add_text_box(slide, "2. Watchlist System", 6, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Separate from favorites",
            "Future viewing intentions",
            "Priority-based organization"
        ], 6, 3, 4, 1.5)
        
        self.add_text_box(slide, "3. Rating & Review System", 11, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "1-5 star rating system",
            "Optional text reviews",
            "User-specific movie ratings"
        ], 11, 3, 4, 1.5)
    
    def add_api_endpoints_slide(self):
        """Slide 9: API Endpoints Architecture"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "API Endpoints Architecture"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "RESTful API Design:", 1, 2, 4, 0.5)
        
        self.add_text_box(slide, "User Management:", 1, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Registration, login, profile management",
            "JWT token refresh",
            "Password change functionality"
        ], 1, 3, 4, 1.5)
        
        self.add_text_box(slide, "Movie Operations:", 6, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Browse, search, filter movies",
            "Detailed movie information",
            "Genre-based categorization"
        ], 6, 3, 4, 1.5)
        
        self.add_text_box(slide, "User Preferences:", 11, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Favorites management",
            "Watchlist operations",
            "Rating and review system"
        ], 11, 3, 4, 1.5)
    
    def add_database_optimization_slide(self):
        """Slide 10: Database Optimization"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Database Optimization"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Performance Strategies:", 1, 2, 4, 0.5)
        
        self.add_text_box(slide, "Indexing Strategy:", 1, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Primary keys on all tables",
            "Foreign key indexes",
            "Composite indexes for junction tables",
            "Unique constraints for data integrity"
        ], 1, 3, 4, 2)
        
        self.add_text_box(slide, "Caching Implementation:", 6, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Redis for session storage",
            "Movie data caching",
            "User preference caching",
            "API response caching"
        ], 6, 3, 4, 2)
        
        self.add_text_box(slide, "Query Optimization:", 11, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Efficient JOIN operations",
            "Pagination for large datasets",
            "Selective field loading",
            "Database connection pooling"
        ], 11, 3, 4, 2)
    
    def add_security_slide(self):
        """Slide 11: Security Implementation"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Security Implementation"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Authentication & Authorization:", 1, 2, 4, 0.5)
        
        self.add_text_box(slide, "JWT Token System:", 1, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Access tokens (60 minutes)",
            "Refresh tokens (24 hours)",
            "Automatic token rotation",
            "Secure token storage"
        ], 1, 3, 4, 2)
        
        self.add_text_box(slide, "Data Protection:", 6, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Password hashing (bcrypt)",
            "Input sanitization",
            "SQL injection prevention",
            "XSS protection"
        ], 6, 3, 4, 2)
        
        self.add_text_box(slide, "API Security:", 11, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Rate limiting",
            "CORS configuration",
            "Request validation",
            "Error handling"
        ], 11, 3, 4, 2)
    
    def add_scalability_slide(self):
        """Slide 12: Scalability Considerations"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Scalability Considerations"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Database Scalability:", 1, 2, 4, 0.5)
        
        self.add_text_box(slide, "Horizontal Scaling:", 1, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Database connection pooling",
            "Read replicas for queries",
            "Sharding strategies for large datasets",
            "Microservices architecture ready"
        ], 1, 3, 4, 2)
        
        self.add_text_box(slide, "Performance Optimization:", 6, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Efficient query patterns",
            "Caching strategies",
            "Pagination implementation",
            "Background task processing"
        ], 6, 3, 4, 2)
        
        self.add_text_box(slide, "Future Enhancements:", 11, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Recommendation algorithms",
            "Machine learning integration",
            "Real-time notifications",
            "Social features"
        ], 11, 3, 4, 2)
    
    def add_deployment_slide(self):
        """Slide 13: Deployment Architecture"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Deployment Architecture"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Production Deployment:", 1, 2, 4, 0.5)
        
        self.add_text_box(slide, "Backend (PythonAnywhere):", 1, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Django application",
            "PostgreSQL database",
            "Redis caching",
            "Static file serving"
        ], 1, 3, 4, 2)
        
        self.add_text_box(slide, "Frontend (Vercel):", 6, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Next.js application",
            "CDN distribution",
            "Automatic deployments",
            "Environment management"
        ], 6, 3, 4, 2)
        
        self.add_text_box(slide, "API Integration:", 11, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "TMDB API for movie data",
            "YouTube API for trailers",
            "Secure API key management"
        ], 11, 3, 4, 1.5)
    
    def add_testing_slide(self):
        """Slide 14: Testing & Quality Assurance"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Testing & Quality Assurance"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Testing Strategy:", 1, 2, 4, 0.5)
        
        self.add_text_box(slide, "Backend Testing:", 1, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Unit tests for models",
            "Integration tests for API endpoints",
            "Authentication testing",
            "Database migration testing"
        ], 1, 3, 4, 2)
        
        self.add_text_box(slide, "Frontend Testing:", 6, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Component testing",
            "User interaction testing",
            "API integration testing",
            "Responsive design testing"
        ], 6, 3, 4, 2)
        
        self.add_text_box(slide, "Quality Metrics:", 11, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Code coverage",
            "Performance benchmarks",
            "Security scanning",
            "Accessibility compliance"
        ], 11, 3, 4, 2)
    
    def add_future_roadmap_slide(self):
        """Slide 15: Future Roadmap"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Future Roadmap"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Planned Enhancements:", 1, 2, 4, 0.5)
        
        self.add_text_box(slide, "Advanced Features:", 1, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Recommendation algorithms",
            "Social sharing capabilities",
            "Advanced search filters",
            "Mobile app development"
        ], 1, 3, 4, 2)
        
        self.add_text_box(slide, "Technical Improvements:", 6, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "GraphQL API implementation",
            "Real-time notifications",
            "Machine learning integration",
            "Microservices architecture"
        ], 6, 3, 4, 2)
        
        self.add_text_box(slide, "User Experience:", 11, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Enhanced UI/UX",
            "Accessibility improvements",
            "Performance optimization",
            "Multi-language support"
        ], 11, 3, 4, 2)
    
    def add_demo_slide(self):
        """Slide 16: Demo & Live Testing"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Demo & Live Testing"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Live Demonstration:", 1, 2, 4, 0.5)
        
        self.add_text_box(slide, "API Testing:", 1, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Swagger documentation access",
            "Endpoint testing",
            "Authentication flow",
            "Data retrieval examples"
        ], 1, 3, 4, 2)
        
        self.add_text_box(slide, "Frontend Features:", 6, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "User registration/login",
            "Movie browsing",
            "Favorites management",
            "Search functionality"
        ], 6, 3, 4, 2)
        
        self.add_text_box(slide, "Database Operations:", 11, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Real-time data queries",
            "User preference management",
            "Performance metrics"
        ], 11, 3, 4, 1.5)
    
    def add_lessons_learned_slide(self):
        """Slide 17: Lessons Learned"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Lessons Learned"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Key Takeaways:", 1, 2, 4, 0.5)
        
        self.add_text_box(slide, "Technical Insights:", 1, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Django ORM best practices",
            "Database design principles",
            "API development patterns",
            "Security implementation"
        ], 1, 3, 4, 2)
        
        self.add_text_box(slide, "Project Management:", 6, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Agile development approach",
            "Version control strategies",
            "Documentation importance",
            "Testing methodologies"
        ], 6, 3, 4, 2)
        
        self.add_text_box(slide, "Personal Growth:", 11, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Full-stack development skills",
            "Database design expertise",
            "API development experience",
            "Deployment knowledge"
        ], 11, 3, 4, 2)
    
    def add_qa_slide(self):
        """Slide 18: Q&A Session"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self.set_dark_background(slide)
        
        title = slide.shapes.title
        title.text = "Q&A Session"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        self.add_text_box(slide, "Questions & Discussion:", 1, 2, 4, 0.5)
        
        self.add_text_box(slide, "Technical Questions:", 1, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Database design decisions",
            "API architecture choices",
            "Security implementation",
            "Performance optimization"
        ], 1, 3, 4, 2)
        
        self.add_text_box(slide, "Project Questions:", 6, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "Development timeline",
            "Challenges encountered",
            "Future enhancements",
            "Learning outcomes"
        ], 6, 3, 4, 2)
        
        self.add_text_box(slide, "Thank You!", 11, 2.5, 4, 0.5)
        self.add_bullet_points(slide, [
            "GitHub Repository: [Link]",
            "Live Demo: [Link]",
            "Documentation: [Link]",
            "Contact: [Your Email]"
        ], 11, 3, 4, 2)
    
    def set_dark_background(self, slide):
        """Set dark background for slide"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 15, 26)
    
    def add_text_box(self, slide, text, left, top, width, height):
        """Add a text box to the slide"""
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    def add_bullet_points(self, slide, points, left, top, width, height):
        """Add bullet points to the slide"""
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.clear()
        
        for i, point in enumerate(points):
            p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.alignment = PP_ALIGN.LEFT
    
    def add_decorative_elements(self, slide):
        """Add decorative elements to the title slide"""
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1), Inches(12), Inches(0.1)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(52, 152, 219)
        rect.line.fill.background()
    
    def create_presentation(self):
        """Create the complete presentation"""
        print("🎬 Creating Movie Recommendation API Presentation...")
        
        # Add all slides
        self.add_title_slide()
        self.add_project_overview_slide()
        self.add_technology_stack_slide()
        self.add_database_architecture_slide()
        self.add_erd_slide()
        self.add_user_management_slide()
        self.add_movie_data_model_slide()
        self.add_user_interaction_slide()
        self.add_api_endpoints_slide()
        self.add_database_optimization_slide()
        self.add_security_slide()
        self.add_scalability_slide()
        self.add_deployment_slide()
        self.add_testing_slide()
        self.add_future_roadmap_slide()
        self.add_demo_slide()
        self.add_lessons_learned_slide()
        self.add_qa_slide()
        
        # Save the presentation
        output_path = "Movie_Recommendation_API_Presentation.pptx"
        self.prs.save(output_path)
        print(f"✅ Presentation saved as: {output_path}")
        print(f"📊 Total slides created: {len(self.prs.slides)}")
        
        return output_path

def main():
    """Main function to create the presentation"""
    try:
        # Create presentation
        presentation = MovieAPIPresentation()
        output_file = presentation.create_presentation()
        
        print("\n🎉 Presentation created successfully!")
        print(f"📁 File location: {os.path.abspath(output_file)}")
        print("\n📋 Next steps:")
        print("1. Open the PowerPoint file")
        print("2. Review and customize content")
        print("3. Add your personal information")
        print("4. Add screenshots or diagrams")
        print("5. Practice your presentation")
        
    except Exception as e:
        print(f"❌ Error creating presentation: {str(e)}")
        print("Make sure you have python-pptx installed:")
        print("pip install python-pptx")

if __name__ == "__main__":
    main() 